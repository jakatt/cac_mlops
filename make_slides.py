"""Generate MLOps architecture PowerPoint — 2 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette Liora ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x12, 0x17, 0x2E)   # fond principal
ORANGE = RGBColor(0xFF, 0x5F, 0x2E)   # accent Liora
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xE8, 0xEA, 0xF0)   # fond boîtes claires
DGREY  = RGBColor(0x2A, 0x30, 0x4D)   # fond boîtes sombres
GREEN  = RGBColor(0x2E, 0xCC, 0x71)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
CYAN   = RGBColor(0x00, 0xBC, 0xD4)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)

# ── Helpers ───────────────────────────────────────────────────────────────────
def px(v): return Inches(v / 96)

def add_bg(slide, color):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, bg, text="", font_size=11,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        border_color=None, border_width=Pt(0), radius=False):
    """Add a rounded-corner text box."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # support multi-paragraph via \n
    paragraphs = text.split("\n")
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    # vertical centering
    from pptx.enum.text import MSO_ANCHOR
    tf.auto_size = None

    return shape

def arrow(slide, x1, y1, x2, y2, color=ORANGE, width=Pt(2)):
    """Add a vertical or horizontal arrow (connector)."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def title_box(slide, text, subtitle=""):
    """Slide header bar."""
    box(slide, 0, 0, 13.33, 0.75, ORANGE, bold=True,
        text=text, font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        box(slide, 0, 0.75, 13.33, 0.35, DGREY,
            text=subtitle, font_size=11, color=LGREY, align=PP_ALIGN.CENTER)

def label(slide, x, y, w, h, text, size=9, color=LGREY, bold=False, align=PP_ALIGN.LEFT):
    box(slide, x, y, w, h, NAVY, text=text, font_size=size,
        color=color, bold=bold, align=align)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — FLOW ANNUEL
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide1(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, NAVY)

    # ── Header ────────────────────────────────────────────────────────────────
    title_box(slide,
              "Flow Annuel — Mise à Jour des Données ONISR",
              "Entraînement : 2021 → 2022 → 2023   ·   Production simulée : données 2024 via simulate_production.py")

    # ── Layout constants ───────────────────────────────────────────────────────
    LX  = 0.18   # left column x
    LW  = 2.80   # left column width
    MX  = 3.20   # main column x
    MW  = 6.70   # main column width
    RX  = 10.10  # right column x  (annotations)
    RW  = 3.05   # right column width
    BH  = 0.52   # box height
    GAP = 0.15   # gap between boxes
    Y0  = 1.22   # first box y

    steps = [
        ("1", "DÉCLENCHEUR",       "Prefect flow planifié\n(annuel) ou manuel",                       DGREY,  ORANGE),
        ("2", "INGESTION",         "data.gouv.fr API · FILENAMES mapping/année\n4 fichiers CSV · ~340 000 lignes",   DGREY,  CYAN),
        ("3", "VALIDATION SCHÉMA", "Pandera · 3 niveaux\nFormat · Colonnes · Qualité",                DGREY,  AMBER),
        ("4", "PREPROCESSING",     "Fusion 4 tables · Feature engineering\n~55 000 lignes × 28 feat", DGREY,  CYAN),
        ("5", "VERSIONING DATA",   "DVC tag data-v{N}\nPush → Scaleway Object Storage",               DGREY,  PURPLE),
        ("6", "ENTRAÎNEMENT",      "RandomForest sur cumul 2021→N\nMLflow run #{N} · params + metrics",DGREY, CYAN),
        ("7", "VALIDATION MODÈLE", "F1 ≥ 0.68 · AUC ≥ 0.75 · Recall ≥ 0.65\nvs modèle en production",DGREY, AMBER),
        ("8", "DÉPLOIEMENT",       "MLflow Registry → Production\nAPI rechargée · rolling update K8s", DGREY, GREEN),
        ("9", "MONITORING",        "Prometheus · Evidently · Grafana\nEvidently : ref=2021-23 · prod=2024 · retrain auto",   DGREY, CYAN),
    ]

    for i, (num, title, desc, bg, accent) in enumerate(steps):
        y = Y0 + i * (BH + GAP)

        # numéro
        box(slide, LX, y, 0.38, BH, accent,
            text=num, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # titre
        box(slide, LX + 0.42, y, LW - 0.42, BH, accent,
            text=title, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

        # description
        box(slide, MX, y, MW, BH, DGREY,
            text=desc, font_size=10, color=LGREY, align=PP_ALIGN.LEFT)

        # flèche vers le bas (sauf dernière)
        if i < len(steps) - 1:
            arrow(slide,
                  LX + LW / 2, y + BH,
                  LX + LW / 2, y + BH + GAP)

    # ── Branche validation (étape 3) ─────────────────────────────────────────
    y3 = Y0 + 2 * (BH + GAP)  # y de l'étape 3

    # ❌ CRITICAL
    box(slide, RX, y3 - 0.02, RW, 0.42, RGBColor(0x4A, 0x10, 0x10),
        text="❌  CRITICAL → STOP\nAlerte équipe · modèle N-1 reste actif",
        font_size=9, color=RED, bold=False)

    # ⚠️ WARNING
    box(slide, RX, y3 + 0.47, RW, 0.38, RGBColor(0x3D, 0x2E, 0x00),
        text="⚠️   WARNING → Log + Continue\nAnomalie enregistrée dans MLflow",
        font_size=9, color=AMBER)

    # ✅ OK
    box(slide, RX, y3 + 0.92, RW, 0.30, RGBColor(0x0A, 0x2E, 0x14),
        text="✅  OK → Pipeline continue",
        font_size=9, color=GREEN)

    # ── Annotation monitoring (étape 9) ─────────────────────────────────────
    y9 = Y0 + 8 * (BH + GAP)  # y de l'étape 9
    box(slide, RX, y9, RW, 0.52, RGBColor(0x0A, 0x25, 0x35),
        text="données 2024 →\nsimulate_production.py → Evidently\nref:2021-23  prod:2024",
        font_size=8, color=CYAN, bold=False)

    # séparateur vertical  avant annotations
    line = slide.shapes.add_shape(1, Inches(RX - 0.08), Inches(y3 - 0.05),
                                  Inches(0.04), Inches(1.30))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()

    # ── Bannière traçabilité ──────────────────────────────────────────────────
    yb = Y0 + len(steps) * (BH + GAP) + 0.05
    box(slide, LX, yb, 13.0, 0.34, DGREY,
        text="  Traçabilité intégrale à chaque run   ·   "
             "Git (code)   ·   DVC (données)   ·   MLflow (modèle + métriques)",
        font_size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # ── Logo Liora (texte) ────────────────────────────────────────────────────
    box(slide, 11.8, 7.15, 1.40, 0.25, NAVY,
        text="Liora — MLOps Accidents", font_size=7, color=DGREY)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ARCHITECTURE & COMPOSANTS
# ═══════════════════════════════════════════════════════════════════════════════
def build_slide2(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg(slide, NAVY)

    title_box(slide,
              "Architecture MLOps — Accidents Routiers",
              "Stack complète · Local → Scaleway · 4 phases")

    # ── Colonnes ──────────────────────────────────────────────────────────────
    # 5 couches empilées : Source · Données · ML · Serving · Monitoring
    # puis split Local / Scaleway
    LY   = 1.18
    FULL = 13.00
    PX   = 0.18

    def layer(y, h, bg, title, tools, title_color=ORANGE, tool_color=WHITE):
        # bande titre à gauche
        box(slide, PX, y, 1.55, h, bg,
            text=title, font_size=10, bold=True,
            color=title_color, align=PP_ALIGN.CENTER)
        # zone outils
        box(slide, PX + 1.60, y, FULL - 1.60, h, DGREY,
            text=tools, font_size=10, color=tool_color, align=PP_ALIGN.LEFT)

    # ── COUCHE 1 : SOURCE ─────────────────────────────────────────────────────
    layer(LY, 0.50, RGBColor(0x1A, 0x35, 0x55),
          "SOURCE",
          "data.gouv.fr / ONISR   ·   Données annuelles (publication juin N+1)   "
          "·   4 fichiers CSV / an   ·   ~340 000 lignes brutes")

    arrow(slide, PX + FULL/2, LY + 0.50, PX + FULL/2, LY + 0.65)

    # ── COUCHE 2 : DONNÉES ────────────────────────────────────────────────────
    y2 = LY + 0.65
    box(slide, PX, y2, 1.55, 0.80, RGBColor(0x2A, 0x1A, 0x55),
        text="DONNÉES", font_size=10, bold=True,
        color=PURPLE, align=PP_ALIGN.CENTER)

    # sous-boîtes
    box(slide, PX+1.60, y2, 3.80, 0.80, DGREY,
        text="Pandera\nValidation schéma\n3 niveaux CRITICAL/WARN/OK",
        font_size=9, color=LGREY)
    box(slide, PX+5.55, y2, 3.80, 0.80, DGREY,
        text="DVC\nVersioning données\ntag data-v{N}",
        font_size=9, color=LGREY)
    box(slide, PX+9.50, y2, 3.85, 0.80, DGREY,
        text="Scaleway Object Storage\nBucket données brutes\n+ preprocessées",
        font_size=9, color=LGREY)

    arrow(slide, PX + FULL/2, y2 + 0.80, PX + FULL/2, y2 + 0.95)

    # ── COUCHE 3 : ML ─────────────────────────────────────────────────────────
    y3 = y2 + 0.95
    box(slide, PX, y3, 1.55, 0.80, RGBColor(0x1A, 0x2E, 0x1A),
        text="ENTRAÎ-\nNEMENT", font_size=10, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)

    box(slide, PX+1.60, y3, 3.80, 0.80, DGREY,
        text="make_dataset.py\nFusion 4 tables · Feature eng.\n~55k lignes × 28 features",
        font_size=9, color=LGREY)
    box(slide, PX+5.55, y3, 3.80, 0.80, DGREY,
        text="train_model.py\nRandomForest (scikit-learn)\nCumul années 2021 → N",
        font_size=9, color=LGREY)
    box(slide, PX+9.50, y3, 3.85, 0.80, DGREY,
        text="MLflow\nTracking · Model Registry\nStaging → Production",
        font_size=9, color=LGREY)

    arrow(slide, PX + FULL/2, y3 + 0.80, PX + FULL/2, y3 + 0.95)

    # ── COUCHE 4 : SERVING ────────────────────────────────────────────────────
    y4 = y3 + 0.95
    box(slide, PX, y4, 1.55, 0.80, RGBColor(0x2E, 0x1A, 0x10),
        text="SERVING", font_size=10, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)

    box(slide, PX+1.60, y4, 3.80, 0.80, DGREY,
        text="NGINX\nReverse proxy · TLS · Auth JWT\nRate limiting 10 req/s",
        font_size=9, color=LGREY)
    box(slide, PX+5.55, y4, 3.80, 0.80, DGREY,
        text="FastAPI + Pydantic\nPOST /predict   GET /health\nGET /metrics",
        font_size=9, color=LGREY)
    box(slide, PX+9.50, y4, 3.85, 0.80, DGREY,
        text="Prefect\nOrchestration des flows\nETL · Train · Retrain",
        font_size=9, color=LGREY)

    arrow(slide, PX + FULL/2, y4 + 0.80, PX + FULL/2, y4 + 0.95)

    # ── COUCHE 5 : MONITORING ─────────────────────────────────────────────────
    y5 = y4 + 0.95
    box(slide, PX, y5, 1.55, 0.80, RGBColor(0x1A, 0x2A, 0x35),
        text="MONI-\nTORING", font_size=10, bold=True,
        color=CYAN, align=PP_ALIGN.CENTER)

    box(slide, PX+1.60, y5, 3.80, 0.80, DGREY,
        text="Prometheus\nMétriques API + pipeline\nLatence · Volume · Erreurs",
        font_size=9, color=LGREY)
    box(slide, PX+5.55, y5, 3.80, 0.80, DGREY,
        text="Evidently\nref: X_train 2021-2023\nprod: données 2024 (simulate_prod.py)",
        font_size=9, color=LGREY)
    box(slide, PX+9.50, y5, 3.85, 0.80, DGREY,
        text="Grafana\nDashboards · Alertes\nDrift CRITICAL → retrain auto",
        font_size=9, color=LGREY)

    # ── SPLIT LOCAL / SCALEWAY ────────────────────────────────────────────────
    y6 = y5 + 0.90
    mid = PX + FULL / 2

    box(slide, PX, y6, FULL / 2 - 0.06, 0.76, RGBColor(0x1A, 0x1F, 0x40),
        text="LOCAL — Docker Compose\n"
             "Services sur localhost · MinIO (S3 local) · "
             "MLflow :5000 · Prefect :4200 · pytest · flake8",
        font_size=9, color=LGREY, border_color=DGREY, border_width=Pt(1))

    box(slide, mid + 0.06, y6, FULL / 2 - 0.06, 0.76, RGBColor(0x1A, 0x1F, 0x40),
        text="SCALEWAY — Kubernetes (Kapsule)\n"
             "N replicas API · Object Storage · Managed DB · "
             "Container Registry · GitHub Actions CI/CD",
        font_size=9, color=LGREY, border_color=DGREY, border_width=Pt(1))

    # séparateur central
    sep = slide.shapes.add_shape(1, Inches(mid + 0.00), Inches(y6),
                                 Inches(0.06), Inches(0.76))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ORANGE
    sep.line.fill.background()

    # label LOCAL / SCALEWAY au-dessus du split
    box(slide, PX,        y6 - 0.22, FULL/2 - 0.06, 0.20, NAVY,
        text="💻  Développement local", font_size=8, color=LGREY, bold=True)
    box(slide, mid + 0.06, y6 - 0.22, FULL/2 - 0.06, 0.20, NAVY,
        text="☁️  Production Scaleway", font_size=8, color=ORANGE, bold=True)

    # ── Légende outils (bas) ──────────────────────────────────────────────────
    box(slide, PX, 7.15, FULL, 0.25, NAVY,
        text="  Git (code)  ·  DVC (données)  ·  MLflow (modèles)  ·  "
             "Pandera (validation)  ·  Prefect (orchestration)  ·  "
             "NGINX (sécurité)  ·  Evidently (drift)  ·  Prometheus/Grafana (monitoring)",
        font_size=8, color=RGBColor(0x88, 0x88, 0xAA), align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs)

    out = "cac_mlops_architecture.pptx"
    prs.save(out)
    print(f"✅  Fichier créé : {out}")

if __name__ == "__main__":
    main()
