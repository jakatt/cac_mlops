#!/usr/bin/env python3
"""
Génère architecture_globale.pptx — diagramme architecture CAC MLOps
Basé sur architecture.md section 6 (état 2026-06-26)

Usage:
    python scripts/generate_architecture_pptx.py
    → génère architecture_globale.pptx dans le répertoire courant
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def rgb(r, g, b):
    return RGBColor(r, g, b)


# ── Palette ────────────────────────────────────────────────────────────────────
NAVY      = rgb(0x1F, 0x3D, 0x6B)
WHITE     = rgb(0xFF, 0xFF, 0xFF)
BLACK     = rgb(0x00, 0x00, 0x00)
BLUE_LT   = rgb(0xDB, 0xE9, 0xF8)   # fond DEV LOCAL
BLUE_HD   = rgb(0x2E, 0x75, 0xB6)   # entête DEV LOCAL
ORG_LT    = rgb(0xFF, 0xF2, 0xCC)   # fond VPS
ORG_HD    = rgb(0xC5, 0x5A, 0x11)   # entête VPS
GREEN_LT  = rgb(0xE2, 0xEF, 0xDA)   # fond KAPSULE
GREEN_HD  = rgb(0x37, 0x86, 0x10)   # entête KAPSULE
PURP_LT   = rgb(0xED, 0xE7, 0xF6)   # fond PREFECT
PURP_HD   = rgb(0x6A, 0x1B, 0x9A)   # entête PREFECT
TEAL_LT   = rgb(0xE0, 0xF7, 0xFA)   # fond MONITORING
TEAL_HD   = rgb(0x00, 0x6D, 0x7E)   # entête MONITORING
GRAD_LT   = rgb(0xFD, 0xE9, 0xD5)   # fond GRADIO
GRAD_HD   = rgb(0xE9, 0x5E, 0x00)   # entête GRADIO
SHAR_LT   = rgb(0xF0, 0xF0, 0xF0)   # fond PARTAGÉ
SHAR_HD   = rgb(0x59, 0x59, 0x59)   # entête PARTAGÉ
TABLE_HDR = rgb(0x2E, 0x75, 0xB6)
TABLE_ROW = rgb(0xF7, 0xFB, 0xFF)
TABLE_ALT = rgb(0xDB, 0xE9, 0xF8)
GRAY_BD   = rgb(0x40, 0x40, 0x40)
RED_TXT   = rgb(0xC0, 0x00, 0x00)
GRAY_TXT  = rgb(0x90, 0x90, 0x90)
BLUE_SUB  = rgb(0xB0, 0xC4, 0xDE)
NAVY_SUB  = rgb(0x8A, 0xA8, 0xD5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, fill=None, border_color=None, border_pt=0.75):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.adjustments[0] = 0.025  # slight rounding
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def _set_text(tf, lines, font_size=7, bold=False, color=BLACK,
              align=PP_ALIGN.LEFT, margin_l=4, margin_t=3):
    tf.word_wrap = True
    tf.margin_left = Pt(margin_l)
    tf.margin_top = Pt(margin_t)
    tf.margin_bottom = Pt(2)
    tf.margin_right = Pt(3)
    if isinstance(lines, str):
        lines = lines.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def header_bar(slide, x, y, w, h, text, fill, font_size=7.5,
               color=WHITE, sub=None):
    shape = _add_rect(slide, x, y, w, h, fill=fill)
    shape.adjustments[0] = 0.02
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(2)
    tf.margin_left = Pt(6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(6)
        r2.font.color.rgb = NAVY_SUB
    return shape


def body_block(slide, x, y, w, h, text, fill, border_color=None,
               font_size=6.5, color=BLACK, border_pt=0.5):
    shape = _add_rect(slide, x, y, w, h, fill=fill,
                      border_color=border_color, border_pt=border_pt)
    shape.adjustments[0] = 0.01
    _set_text(shape.text_frame, text, font_size=font_size, color=color)
    return shape


def section_header(slide, x, y, w, h, text, fill, font_size=6.5, color=WHITE):
    shape = _add_rect(slide, x, y, w, h, fill=fill)
    shape.adjustments[0] = 0.01
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(2)
    tf.margin_left = Pt(5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = color
    return shape


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── TITLE BAR ──────────────────────────────────────────────────────────────
    title_h = 1.15
    title = _add_rect(slide, 0, 0, 33.87, title_h, fill=NAVY)
    title.adjustments[0] = 0
    tf = title.text_frame
    tf.margin_top = Pt(5)
    tf.margin_left = Pt(8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "ARCHITECTURE GLOBALE — CAC MLOPS"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = ("VPS Scaleway DEV1-XL  ·  Docker 13 conteneurs  ·  "
               "Prefect 11 flows  ·  7 workflows GitHub Actions")
    r2.font.size = Pt(7)
    r2.font.color.rgb = BLUE_SUB

    # ── LAYOUT ─────────────────────────────────────────────────────────────────
    Y0  = title_h + 0.1
    H   = 19.05 - Y0 - 0.1   # 17.8 cm available

    x_dev = 0.1
    w_dev = 5.2
    x_vps = x_dev + w_dev + 0.1
    w_vps = 20.0
    x_kap = x_vps + w_vps + 0.1
    w_kap = 33.87 - x_kap - 0.12   # ~8.35 cm

    # ── COL 1 — DEV LOCAL ──────────────────────────────────────────────────────
    _add_rect(slide, x_dev, Y0, w_dev, H,
              fill=BLUE_LT, border_color=BLUE_HD, border_pt=1.5)

    header_bar(slide, x_dev, Y0, w_dev, 1.1,
               "DEV LOCAL", fill=BLUE_HD,
               sub="Mac développeur")

    body_block(slide, x_dev + 0.12, Y0 + 1.25, w_dev - 0.24, 5.0,
               "docker-compose.yml\n"
               "(même stack que VPS)\n"
               "ports → 127.0.0.1\n"
               "volumes → ./\n\n"
               "Outils CLI\n"
               "──────────────────\n"
               "git · dvc · pytest\n"
               "flake8 · kubectl\n\n"
               "DVC pull\n"
               "→ data/  (S3 Scaleway)\n\n"
               "MLFLOW_TRACKING_URI\n"
               "100.117.99.62:5001\n"
               "(via Tailscale)\n\n"
               "MLFLOW_S3_ENDPOINT\n"
               "100.117.99.62:9000",
               fill=BLUE_LT, font_size=7, color=NAVY)

    section_header(slide, x_dev + 0.2, Y0 + 6.5, w_dev - 0.4, 0.55,
                   "↔  git push / dvc push", fill=NAVY)

    body_block(slide, x_dev + 0.12, Y0 + 7.2, w_dev - 0.24, 4.5,
               "Cycle quotidien DS\n"
               "──────────────────\n"
               "git pull && dvc pull\n"
               "→ dev + expériences\n"
               "→ git push + dvc push\n"
               "→ PR vers main\n"
               "→ deploy automatique\n\n"
               "Cycle annuel ONISR\n"
               "──────────────────\n"
               "→ etl flow (Prefect)\n"
               "→ train flow\n"
               "→ drift check",
               fill=BLUE_LT, font_size=7, color=GRAY_BD)

    # ── COL 2 — VPS SCALEWAY ───────────────────────────────────────────────────
    _add_rect(slide, x_vps, Y0, w_vps, H,
              fill=ORG_LT, border_color=ORG_HD, border_pt=1.5)

    header_bar(slide, x_vps, Y0, w_vps, 1.1,
               "VPS SCALEWAY — DEV1-XL  (Scaleway fr-par-2)", fill=ORG_HD,
               sub=("IP publique : 51.159.187.132   ·   IP Tailscale : 100.117.99.62   ·   "
                    "/ = 20 GB NVMe   ·   /data = 80 GB block storage"))

    # VPS split: left 9.0cm containers | gap 0.15 | right remainder
    x_vl  = x_vps + 0.15
    w_vl  = 9.0
    x_vr  = x_vl + w_vl + 0.15
    w_vr  = w_vps - w_vl - 0.45
    y_sub = Y0 + 1.25

    # ── VPS LEFT — CONTAINERS TABLE ─────────────────────────────────
    section_header(slide, x_vl, y_sub, w_vl, 0.55,
                   "CONTAINERS DOCKER  (13 : 12 permanents + minio-init EXIT)",
                   fill=TABLE_HDR)

    CONTAINERS = [
        ("postgresql",     "5432",        "interne Docker"),
        ("minio",          "9000 / 9001", "Tailscale"),
        ("minio-init",     "—",           "EXIT après init"),
        ("mlflow",         "5001",        "Tailscale"),
        ("api",            "8080 / 8000", "Tailscale / Prom."),
        ("nginx",          "8090",        "PUBLIC  0.0.0.0"),
        ("prefect-server", "4200",        "Tailscale"),
        ("prefect-worker", "—",           "process pool"),
        ("gradio",         "7860",        "Tailscale"),
        ("node-exporter",  "9100",        "interne Docker"),
        ("nginx-exporter", "9113",        "interne Docker"),
        ("prometheus",     "9090",        "Tailscale"),
        ("grafana",        "3000",        "Tailscale"),
    ]

    tbl_y = y_sub + 0.6
    tbl_h = 9.6
    n_rows = len(CONTAINERS) + 1   # header row + 13 data rows
    tbl_shape = slide.shapes.add_table(
        n_rows, 3,
        Cm(x_vl), Cm(tbl_y), Cm(w_vl), Cm(tbl_h)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Cm(3.1)
    tbl.columns[1].width = Cm(2.3)
    tbl.columns[2].width = Cm(3.6)

    for j, hdr in enumerate(("Conteneur", "Port hôte", "Accès")):
        cell = tbl.cell(0, j)
        cell.text = hdr
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(7)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR

    for i, (name, port, access) in enumerate(CONTAINERS):
        row_fill = TABLE_ALT if i % 2 == 0 else TABLE_ROW
        is_public = "PUBLIC" in access
        is_exit   = "EXIT"   in access
        for j, val in enumerate((name, port, access)):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(6.5)
            run.font.bold = (j == 0) or (is_public and j == 2)
            run.font.color.rgb = (
                RED_TXT  if is_public and j == 2 else
                GRAY_TXT if is_exit else
                BLACK
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill

    leg_y = tbl_y + tbl_h + 0.12
    body_block(slide, x_vl, leg_y, w_vl, H - (leg_y - Y0) - 0.15,
               "Niveaux d'accès\n"
               "PUBLIC      → 51.159.187.132 (internet ouvert)\n"
               "Tailscale   → 100.117.99.62  (équipe VPN uniquement)\n"
               "interne     → réseau Docker (non routable hors Docker)\n"
               "process pool→ aucun port exposé (exécute les flows)\n"
               "EXIT        → one-shot : s'arrête après l'initialisation",
               fill=ORG_LT, font_size=6.5, color=GRAY_BD)

    # ── VPS RIGHT — PREFECT / MONITORING / GRADIO ───────────────────
    y_cur = y_sub

    # PREFECT
    section_header(slide, x_vr, y_cur, w_vr, 0.55,
                   "ORCHESTRATION — PREFECT  (11 deployments)",
                   fill=PURP_HD)
    y_cur += 0.6
    pref_h = 6.0
    body_block(slide, x_vr, y_cur, w_vr, pref_h,
               "prefect-server   :4200  (Tailscale)\n"
               "prefect-worker   image api + kubectl + scw + docker CLI\n"
               "pool             default-process-pool (type: process)\n\n"
               "ML / ETL flows\n"
               "  etl            download data.gouv.fr + preprocessing\n"
               "  train          benchmark RF/XGBoost/LGBM + promote\n"
               "  retrain-annual réentraînement annuel (1 algo)\n"
               "  drift-check    drift mensuel Evidently\n"
               "  full-retrain   tous les cycles depuis zéro\n"
               "  reset          vide predictions + rapports drift\n"
               "  check-new-data détection données ONISR (lundi 8h)\n\n"
               "Infra / Ops flows\n"
               "  kapsule-up     provision cluster K8s + upload modèle S3\n"
               "  kapsule-down   déprovision cluster K8s\n"
               "  test-api       tests end-to-end (JWT, /predict, 429)\n"
               "  diag           diagnostic VPS (disk, docker, network)",
               fill=PURP_LT, border_color=PURP_HD)
    y_cur += pref_h + 0.12

    # MONITORING
    section_header(slide, x_vr, y_cur, w_vr, 0.55,
                   "MONITORING", fill=TEAL_HD)
    y_cur += 0.6
    mon_h = 4.5
    body_block(slide, x_vr, y_cur, w_vr, mon_h,
               "Prometheus  :9090  (Tailscale)\n"
               "Scrape targets\n"
               "  api:8000/metrics      req, latence, drift, prédictions\n"
               "  node-exporter:9100    CPU / RAM / disk\n"
               "  nginx-exporter:9113   connexions nginx\n\n"
               "Grafana  :3000  (Tailscale)\n"
               "  api-performance.json  latence p50/p95/p99, taux 5xx\n"
               "  model-drift.json      drift_share, features driftées\n\n"
               "4 alertes email\n"
               "  brute-force 401 (>20/5min)  ·  DDoS 429 (>50/5min)\n"
               "  RAM < 10%  ·  Disk /data < 15%",
               fill=TEAL_LT, border_color=TEAL_HD)
    y_cur += mon_h + 0.12

    # GRADIO
    section_header(slide, x_vr, y_cur, w_vr, 0.55,
                   "COCKPIT — GRADIO  :7860  (Tailscale)", fill=GRAD_HD)
    y_cur += 0.6
    grad_h = H - (y_cur - Y0) - 0.1
    body_block(slide, x_vr, y_cur, w_vr, grad_h,
               "7 onglets\n"
               "  1. What-If        scénarios météo/nuit/alcool/vitesse\n"
               "  2. Points Noirs   density_mapbox France, filtres\n"
               "  3. Drift          rapports Evidently (iframe HTML)\n"
               "  4. Modèles        MLflow Registry, métriques, @Production\n"
               "  5. Pipeline       déclenchement flows Prefect + runs récents\n"
               "  6. Healthcheck    état HTTP tous services VPS + Kapsule\n"
               "  7. Infra          URLs Tailscale + API publique + IPs Kapsule",
               fill=GRAD_LT, border_color=GRAD_HD)

    # ── COL 3 — KAPSULE + PARTAGÉ ──────────────────────────────────────────────
    _add_rect(slide, x_kap, Y0, w_kap, H,
              fill=GREEN_LT, border_color=GREEN_HD, border_pt=1.5)

    # KAPSULE
    header_bar(slide, x_kap, Y0, w_kap, 1.1,
               "KAPSULE K8s  (on-demand)", fill=GREEN_HD,
               sub="Scaleway fr-par  ·  cluster: cac-mlops  ·  K8s 1.35.3")

    kap_content_h = 9.6
    body_block(slide, x_kap + 0.12, Y0 + 1.25, w_kap - 0.24, kap_content_h,
               "Deployments  (namespace: cac-mlops)\n"
               "────────────────────────────────────────\n"
               "api           HPA min 1 → max 8 pods\n"
               "              initContainer fetch S3 → /app/model/\n"
               "mlflow        SQLite emptyDir + s3://cac-mlops-data\n"
               "prefect-server\n"
               "prefect-worker\n"
               "prometheus    scrape api:8000/metrics\n"
               "grafana       ConfigMaps provisionnés\n\n"
               "LoadBalancers LB-S\n"
               "────────────────────────────────────────\n"
               "nginx    :80   → API publique (rate-limit)\n"
               "prefect  :4200 → UI Prefect\n"
               "grafana  :3000 → dashboards\n"
               "mlflow   port-forward uniquement\n\n"
               "HPA api\n"
               "CPU 70% / RAM 80%  ·  min 1 → max 8 pods\n\n"
               "Secrets K8s\n"
               "────────────────────────────────────────\n"
               "s3-creds  : AWS_ACCESS_KEY_ID\n"
               "            AWS_SECRET_ACCESS_KEY\n"
               "app-creds : JWT_SECRET_KEY\n"
               "            API_USERNAME / API_PASSWORD\n\n"
               "État cluster\n"
               "state/kapsule_ips (VPS)\n"
               "→ écrit par flow kapsule-up\n"
               "→ lu par Gradio onglet Infra",
               fill=GREEN_LT, border_color=GREEN_HD, font_size=6.5)

    # PARTAGÉ
    y_shr = Y0 + 1.25 + kap_content_h + 0.15
    shr_h = H - (y_shr - Y0) - 0.1
    section_header(slide, x_kap, y_shr, w_kap, 0.55,
                   "PARTAGÉ", fill=SHAR_HD)
    body_block(slide, x_kap + 0.12, y_shr + 0.6, w_kap - 0.24, shr_h - 0.65,
               "GitHub  jakatt/cac_mlops\n"
               "7 workflows CI/CD\n"
               "  ci · deploy · train\n"
               "  promote · drift\n"
               "  benchmark · cleanup\n\n"
               "GHCR  ghcr.io/jakatt/\n"
               "  cac-mlops-api:latest\n"
               "  cac-mlops-mlflow:latest\n"
               "  cac-mlops-gradio:latest\n\n"
               "Scaleway Object Storage\n"
               "s3://cac-mlops-data\n"
               "  dvc/       → données DVC\n"
               "  k8s-model/ → modèle K8s init\n"
               "  mlflow-k8s/→ artefacts K8s\n\n"
               "MinIO (VPS)\n"
               "→ artefacts MLflow local\n\n"
               "data.gouv.fr  (ONISR)\n"
               "accidents France 2021→2024",
               fill=SHAR_LT, border_color=SHAR_HD, font_size=6.5)

    # ── Save ───────────────────────────────────────────────────────────────────
    out = "architecture_globale.pptx"
    prs.save(out)
    print(f"✓  {out}  généré")
    print("   Ouvrir dans PowerPoint → toutes les formes sont éditables nativement")


if __name__ == "__main__":
    main()
